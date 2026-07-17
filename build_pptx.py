"""
Génère la présentation PowerPoint professionnelle pour la Direction Générale,
selon la charte graphique Afri-Farmers Market (voir brand.py).
"""
import sys
sys.path.insert(0, ".")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import brand
from src_pipeline import load_clean
from src_analysis import performance_globale, analyse_categorie, analyse_ville, analyse_marketing, analyse_clients
from src_viz_brand import generate_all

RGB = lambda h: RGBColor(*brand.hex_to_rgb(h))

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color_hex):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGB(color_hex)


def add_rect(slide, x, y, w, h, color_hex, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGB(color_hex)
    if line:
        shape.line.color.rgb = RGB(color_hex)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=brand.INK, bold=False, align=PP_ALIGN.LEFT,
             font=brand.FONT_BODY, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = font
            run.font.color.rgb = RGB(color)
    return box


def add_bullets(slide, x, y, w, h, items, size=16, color=brand.INK, bullet_color=brand.ACCENT, spacing=10):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"●  {item}"
        run.font.size = Pt(size)
        run.font.name = brand.FONT_BODY
        run.font.color.rgb = RGB(color)
    return box


def add_logo(slide, x, y, size, mono_white_bg=None):
    slide.shapes.add_picture(brand.LOGO_PATH, x, y, height=size)


def add_page_chrome(slide, section, page_num, dark=False):
    fg = brand.WHITE if dark else brand.GREY
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.35), section, size=10,
              color=fg, italic=True)
    add_text(slide, Inches(11.6), Inches(7.05), Inches(0.55), Inches(0.35), str(page_num), size=10,
              color=fg, align=PP_ALIGN.RIGHT)
    add_logo(slide, Inches(12.65), Inches(7.02), Inches(0.35))


def add_header(slide, kicker, title):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.35), brand.PRIMARY)
    add_text(slide, Inches(0.6), Inches(0.18), Inches(9), Inches(0.35), kicker.upper(), size=13,
              color=brand.ACCENT_LIGHT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.5), Inches(10.5), Inches(0.7), title, size=26, bold=True,
              color=brand.WHITE)
    add_logo(slide, Inches(12.35), Inches(0.28), Inches(0.8))


def content_slide(kicker, title, section, page_num):
    slide = add_slide()
    set_background(slide, brand.WHITE)
    add_header(slide, kicker, title)
    add_page_chrome(slide, section, page_num)
    return slide


# ============================================================ DATA
df = load_clean()
perf = performance_globale(df)
cat_agg = analyse_categorie(df)
ville_agg = analyse_ville(df)
mkt_agg = analyse_marketing(df)
clients_info = analyse_clients(df)
charts = generate_all(df)

top_cat = cat_agg.sort_values("ca", ascending=False).iloc[0]
worst_return_cat = cat_agg.sort_values("taux_retour", ascending=False).iloc[0]
top_ville = ville_agg.sort_values("ca", ascending=False).iloc[0]
worst_annulation_ville = ville_agg.sort_values("taux_annulation", ascending=False).iloc[0]
best_roi_canal = mkt_agg.sort_values("roi", ascending=False).iloc[0]
worst_roi_canal = mkt_agg.sort_values("roi", ascending=True).iloc[0]
biggest_budget_canal = mkt_agg.sort_values("cout_marketing", ascending=False).iloc[0]
second_worst_annulation_ville = ville_agg.sort_values("taux_annulation", ascending=False).iloc[1]

PAGE = {"n": 0}


def next_page():
    PAGE["n"] += 1
    return PAGE["n"]


# ============================================================ 1. TITLE SLIDE
slide = add_slide()
set_background(slide, brand.PRIMARY)
add_rect(slide, 0, Inches(6.9), SLIDE_W, Inches(0.1), brand.ACCENT)
add_logo(slide, Inches(5.67), Inches(0.7), Inches(2.0))
add_text(slide, Inches(1), Inches(3.0), Inches(11.33), Inches(1.0), "Analyse Stratégique AfriMarket",
          size=40, bold=True, color=brand.WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.95), Inches(11.33), Inches(0.6),
          "6 mois d'activité commerciale — Présentation à la Direction Générale",
          size=18, color=brand.ACCENT_LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.5), Inches(11.33), Inches(0.4), "Direction Data & Analyse — Juillet 2026",
          size=13, color=brand.WHITE, align=PP_ALIGN.CENTER, italic=True)
next_page()

# ============================================================ 2. SOMMAIRE
slide = add_slide()
set_background(slide, brand.WHITE)
add_rect(slide, 0, 0, Inches(4.2), SLIDE_H, brand.PRIMARY)
add_logo(slide, Inches(1.6), Inches(0.6), Inches(1.0))
add_text(slide, Inches(0.5), Inches(2.0), Inches(3.2), Inches(1.0), "Sommaire", size=30, bold=True, color=brand.WHITE)
items = [
    "Contexte & objectifs",
    "Méthodologie & qualité des données",
    "Performance globale",
    "Analyse par catégorie",
    "Analyse géographique",
    "Analyse marketing",
    "Analyse clients",
    "Recommandations stratégiques",
    "Plan d'action & conclusion",
]
box = slide.shapes.add_textbox(Inches(4.7), Inches(0.7), Inches(8.0), Inches(6))
tf = box.text_frame
tf.word_wrap = True
for i, item in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(18)
    r1 = p.add_run()
    r1.text = f"{i+1:02d}   "
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = RGB(brand.ACCENT_DARK)
    r1.font.name = brand.FONT_BODY
    r2 = p.add_run()
    r2.text = item
    r2.font.size = Pt(20)
    r2.font.color.rgb = RGB(brand.INK)
    r2.font.name = brand.FONT_BODY
next_page()

# ============================================================ 3. CONTEXTE
slide = content_slide("Contexte", "Une entreprise en croissance, des signaux à clarifier", "Contexte & objectifs", next_page())
add_bullets(slide, Inches(0.7), Inches(1.7), Inches(11.8), Inches(2.2), [
    "AfriMarket est un e-commerce panafricain actif dans 8 villes d'Afrique francophone, sur 4 catégories : "
    "Électronique, Mode, Beauté, Maison.",
    "La direction observe des variations de chiffre d'affaires, un taux de retour préoccupant, des dépenses "
    "marketing élevées et des écarts de performance entre villes.",
], size=17, spacing=16)
add_rect(slide, Inches(0.7), Inches(4.1), Inches(11.9), Inches(2.3), brand.LIGHT_BG)
add_text(slide, Inches(1.0), Inches(4.3), Inches(11.2), Inches(0.4), "Objectif de l'analyse", size=16, bold=True,
          color=brand.PRIMARY)
add_text(slide, Inches(1.0), Inches(4.8), Inches(11.2), Inches(1.5),
          "Produire une analyse stratégique complète, fondée sur 6 mois de données commerciales, pour permettre "
          "à la direction de prioriser ses décisions : où investir, quoi corriger, et comment améliorer la "
          "rentabilité et la fidélisation client.", size=15, color=brand.INK, line_spacing=1.25)

# ============================================================ 4. METHODOLOGIE
slide = content_slide("Méthodologie", "Des données réelles, avec leurs imperfections", "Méthodologie", next_page())
add_text(slide, Inches(0.7), Inches(1.6), Inches(11.9), Inches(0.4),
          "10 100 commandes brutes auditées → 10 000 commandes propres et fiables (df_clean)", size=17, bold=True,
          color=brand.PRIMARY)

issues = [
    ("100", "commandes dupliquées", "supprimées (id_commande unique)"),
    ("632", "prix unitaires négatifs", "≈97% figés à -50 (code d'erreur) → imputés par la médiane du produit"),
    ("608", "quantités nulles", "incohérentes business → imputées par la médiane de la catégorie"),
    ("614", "remises négatives", "constante -0.10 (inversion de signe) → valeur absolue"),
]
x0 = Inches(0.7)
w = Inches(2.85)
gap = Inches(0.15)
for i, (num, label, detail) in enumerate(issues):
    x = Emu(x0 + i * (w + gap))
    add_rect(slide, x, Inches(2.3), w, Inches(2.7), brand.LIGHT_BG)
    add_text(slide, x, Inches(2.45), w, Inches(0.7), num, size=32, bold=True, color=brand.PRIMARY, align=PP_ALIGN.CENTER)
    add_text(slide, Emu(x + Inches(0.15)), Inches(3.15), Emu(w - Inches(0.3)), Inches(0.5), label, size=13, bold=True,
              color=brand.ACCENT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, Emu(x + Inches(0.15)), Inches(3.65), Emu(w - Inches(0.3)), Inches(1.2), detail, size=11,
              color=brand.GREY, align=PP_ALIGN.CENTER, line_spacing=1.15)

add_bullets(slide, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.4), [
    "Incohérences orthographiques et de casse harmonisées (villes, catégories, statuts).",
    "Hypothèse documentée : marge brute = 45% du CA (coût de revient non fourni dans les données).",
], size=14, spacing=8)

# ============================================================ 5. KPIs GLOBAUX
slide = content_slide("Performance globale", "Une activité saine, un signal de qualité à traiter", "Performance globale", next_page())
kpis = [
    (f"{perf['ca_total']:,.0f}".replace(",", " "), "Chiffre d'affaires total"),
    (f"{perf['profit_total']:,.0f}".replace(",", " "), "Profit net estimé"),
    (f"{perf['panier_moyen']:,.0f}".replace(",", " "), "Panier moyen"),
    (f"{perf['taux_annulation']:.1%}", "Taux d'annulation"),
    (f"{perf['taux_retour']:.1%}", "Taux de retour"),
    (f"{perf['nb_clients']:,}".replace(",", " "), "Clients actifs"),
]
x0 = Inches(0.6)
w = Inches(1.98)
gap = Inches(0.08)
for i, (val, label) in enumerate(kpis):
    x = Emu(x0 + i * (w + gap))
    add_rect(slide, x, Inches(1.9), w, Inches(2.1), brand.PRIMARY)
    add_text(slide, x, Inches(2.15), w, Inches(0.8), val, size=22, bold=True, color=brand.WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Emu(x + Inches(0.1)), Inches(3.1), Emu(w - Inches(0.2)), Inches(0.8), label, size=11.5,
              color=brand.ACCENT_LIGHT, align=PP_ALIGN.CENTER, line_spacing=1.1)

add_bullets(slide, Inches(0.7), Inches(4.5), Inches(11.9), Inches(2.2), [
    "L'activité est globalement saine : le taux d'annulation reste très faible (1,9%).",
    "Le taux de retour (8,2%) mérite attention : il est très concentré sur une seule catégorie (voir slide suivante).",
], size=16, spacing=14)

# ============================================================ 6. CATEGORIE - CA
slide = content_slide("Analyse par catégorie", f"{top_cat['categorie']} : le moteur économique d'AfriMarket",
                        "Analyse par catégorie", next_page())
slide.shapes.add_picture(charts["ca_categorie"], Inches(0.6), Inches(1.6), height=Inches(4.3))
add_bullets(slide, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4), [
    f"{top_cat['categorie']} génère {top_cat['ca']/perf['ca_total']:.0%} du chiffre d'affaires total.",
    f"Elle porte l'essentiel du profit net ({top_cat['profit']:,.0f}".replace(",", " ") + " sur "
    f"{perf['profit_total']:,.0f}".replace(",", " ") + " au global).",
    "C'est la catégorie à prioriser en stock, visibilité et négociation fournisseur.",
], size=16, spacing=16)
next_page()

# ============================================================ 7. CATEGORIE - RETOUR
slide = content_slide("Analyse par catégorie", "Un taux de retour à traiter en priorité", "Analyse par catégorie", PAGE["n"])
slide.shapes.add_picture(charts["retour_categorie"], Inches(0.6), Inches(1.6), height=Inches(4.3))
add_bullets(slide, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4), [
    f"{worst_return_cat['categorie']} affiche un taux de retour de {worst_return_cat['taux_retour']:.1%}, "
    f"2 à 5 fois supérieur aux autres catégories.",
    "Chaque point de retour y coûte proportionnellement plus cher (logistique inverse, remboursement, image).",
    "Un plan qualité ciblé (contrôle fournisseur, fiches produit, politique de retour) est nécessaire.",
], size=16, spacing=16)

# ============================================================ 8. GEOGRAPHIE - CA
slide = content_slide("Analyse géographique", f"{top_ville['ville']} : la ville la plus performante", "Analyse géographique", next_page())
slide.shapes.add_picture(charts["ca_ville"], Inches(0.5), Inches(1.6), height=Inches(4.4))
add_bullets(slide, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4), [
    f"{top_ville['ville']} combine le meilleur CA ({top_ville['ca']:,.0f}".replace(",", " ") + ") et un taux "
    f"d'annulation quasi nul ({top_ville['taux_annulation']:.1%}).",
    "C'est la ville à privilégier pour un investissement additionnel (stock, marketing local).",
], size=16, spacing=16)
next_page()

# ============================================================ 9. GEOGRAPHIE - ANNULATION
slide = content_slide("Analyse géographique", f"{worst_annulation_ville['ville']} : un signal opérationnel à investiguer",
                        "Analyse géographique", PAGE["n"])
slide.shapes.add_picture(charts["annulation_ville"], Inches(0.5), Inches(1.6), height=Inches(4.4))
add_bullets(slide, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4), [
    f"{worst_annulation_ville['ville']} affiche un taux d'annulation de {worst_annulation_ville['taux_annulation']:.1%}, "
    "contre quasiment 0% dans les autres villes.",
    "Ce signal localisé (paiement, logistique, expérience client) doit être diagnostiqué avant tout "
    "investissement marketing supplémentaire dans cette ville.",
], size=16, spacing=16)

# ============================================================ 10. MARKETING
slide = content_slide("Analyse marketing", "Un budget à rééquilibrer vers l'efficacité", "Analyse marketing", next_page())
slide.shapes.add_picture(charts["roi_canal"], Inches(0.6), Inches(1.6), height=Inches(4.3))
add_bullets(slide, Inches(7.3), Inches(2.0), Inches(5.4), Inches(4), [
    f"{best_roi_canal['canal_marketing']} affiche le meilleur ROI (x{best_roi_canal['roi']:.0f}) pour un coût "
    "marginal très faible.",
    f"{biggest_budget_canal['canal_marketing']} capte le budget le plus élevé mais un ROI modeste "
    f"(x{biggest_budget_canal['roi']:.0f}) : le plus gros gisement d'optimisation en valeur absolue.",
    f"{worst_roi_canal['canal_marketing']} affiche le ROI le plus faible en proportion (x{worst_roi_canal['roi']:.0f}), sur un budget plus restreint.",
    "Recommandation : réallouer progressivement le budget vers les canaux les plus efficaces (test A/B).",
], size=15, spacing=13)
add_text(slide, Inches(0.6), Inches(6.0), Inches(6), Inches(0.4), "ROI = (CA − coût marketing) / coût marketing",
          size=11, italic=True, color=brand.GREY)

# ============================================================ 11. CLIENTS
slide = content_slide("Analyse clients", "Une base fidèle, mais concentrée", "Analyse clients", next_page())
slide.shapes.add_picture(charts["pareto"], Inches(0.5), Inches(1.6), height=Inches(4.2))
slide.shapes.add_picture(charts["segmentation"], Inches(7.4), Inches(1.7), height=Inches(4.0))
add_text(slide, Inches(0.6), Inches(5.9), Inches(11.9), Inches(1.2),
          f"{clients_info['pct_recurrents']:.0%} des {clients_info['nb_clients']:,}".replace(",", " ") +
          f" clients actifs sont récurrents, et les 20% les plus rentables génèrent {charts['part20']}% du CA total.",
          size=16, color=brand.INK, bold=True)

# ============================================================ 12. RECOMMANDATIONS
slide = add_slide()
set_background(slide, brand.WHITE)
next_page()
add_header(slide, "Recommandations", "5 recommandations stratégiques")
recos = [
    "Sécuriser Électronique (75% du CA) tout en lançant un plan qualité pour réduire son taux de retour.",
    f"Investir davantage à {top_ville['ville']} ; auditer {worst_annulation_ville['ville']} avant tout budget supplémentaire.",
    f"Réallouer une partie du budget {biggest_budget_canal['canal_marketing']} (x{biggest_budget_canal['roi']:.0f}) vers {best_roi_canal['canal_marketing']} (x{best_roi_canal['roi']:.0f}) ; revoir aussi {worst_roi_canal['canal_marketing']}, dont le ROI (x{worst_roi_canal['roi']:.0f}) est le plus faible en proportion.",
    "Lancer un programme de fidélisation dédié aux clients VIP pour sécuriser leur rente.",
    "Mettre en place une relance automatisée pour les clients occasionnels (une seule commande).",
]
y = Inches(1.75)
for i, reco in enumerate(recos):
    add_rect(slide, Inches(0.7), y, Inches(0.7), Inches(0.7), brand.PRIMARY)
    add_text(slide, Inches(0.7), y, Inches(0.7), Inches(0.7), str(i + 1), size=26, bold=True, color=brand.WHITE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(1.65), y, Inches(10.8), Inches(0.9), reco, size=16, color=brand.INK,
              anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    y = Emu(y + Inches(1.02))
add_page_chrome(slide, "Recommandations", PAGE["n"])

# ============================================================ 13. PLAN D'ACTION
slide = content_slide("Plan d'action", "Une feuille de route en 3 horizons", "Plan d'action", next_page())
horizons = [
    ("0-30 jours", ["Audit opérationnel de Douala (coût faible, risque élevé à ne rien faire)",
                     "Test A/B de réallocation budget Instagram → Email"]),
    ("30-90 jours", ["Lancement du programme de fidélisation VIP",
                       "Séquence de relance des clients occasionnels"]),
    ("2-3 trimestres", ["Plan qualité Électronique (fournisseurs, fiches produit, retours)",
                          "Renforcement de l'investissement à Kinshasa"]),
]
x0 = Inches(0.7)
w = Inches(3.85)
gap = Inches(0.25)
for i, (h, items) in enumerate(horizons):
    x = Emu(x0 + i * (w + gap))
    add_rect(slide, x, Inches(1.8), w, Inches(0.6), brand.ACCENT_DARK)
    add_text(slide, x, Inches(1.8), w, Inches(0.6), h, size=17, bold=True, color=brand.WHITE, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, x, Inches(2.5), w, Inches(3.6), brand.LIGHT_BG)
    add_bullets(slide, Emu(x + Inches(0.2)), Inches(2.7), Emu(w - Inches(0.4)), Inches(3.2), items, size=13.5,
                 spacing=14)

# ============================================================ 14. CONCLUSION
slide = content_slide("Conclusion", "Une performance rentable, mais concentrée et donc fragile", "Conclusion", next_page())
add_text(slide, Inches(0.7), Inches(1.8), Inches(11.9), Inches(3.5),
          "AfriMarket dispose d'un modèle rentable et d'une base client majoritairement fidèle. Mais sa "
          "performance repose sur des équilibres fragiles : une catégorie (Électronique) porte l'essentiel du "
          "chiffre d'affaires, une ville (Douala) présente une anomalie opérationnelle non expliquée, et un canal "
          "marketing (Instagram Ads) absorbe un budget disproportionné à son efficacité réelle.\n\n"
          "Les prochaines 90 jours doivent prioriser les actions à fort impact et faible coût — audit Douala, "
          "réallocation budgétaire, programme VIP — avant d'engager les chantiers plus structurants.",
          size=18, color=brand.INK, line_spacing=1.4)

# ============================================================ 15. CLOSING
slide = add_slide()
set_background(slide, brand.PRIMARY)
add_logo(slide, Inches(5.67), Inches(1.6), Inches(1.8))
add_text(slide, Inches(1), Inches(3.7), Inches(11.33), Inches(0.8), "Merci de votre attention",
          size=34, bold=True, color=brand.WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(0.5), "Questions & discussion",
          size=18, color=brand.ACCENT_LIGHT, align=PP_ALIGN.CENTER)

prs.save("reports/AfriMarket_Presentation_Direction.pptx")
print("PPTX généré :", len(prs.slides.__iter__.__self__._sldIdLst), "diapositives")
